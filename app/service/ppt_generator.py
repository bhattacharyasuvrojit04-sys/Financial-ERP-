from pptx import Presentation

def built_pitch_deck(deck_data, output_file):
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])

    slide.shapes.title.text = "Executive Summary"

    slide.placeholders[1].text = \
        deck_data["executive_summary"]
    
    #slide 2:

    slide = prs.slides.add_slides(prs.slide_layouts[1])

    slide.shapes.title.text = "Company Overview"

    slide.placeholders[1].text = \
        deck_data["company_overview"]
    
    #slide 3:

    slide = prs.slides.add_slides(prs.slide_layouts[1])

    slide.shapes.title.text = "Financial Highlights"

    slide.placeholders[1].text = \
        deck_data["financial_highlights"]
    
    #slide 4:
    slide = prs.slides.add_slides(prs.slide_layouts[1])

    slide.shapes.title.text = "SWOT Analysis"

    slide.placeholders[1].text = f"""
        Strengths:
        {deck_data["swot_strengths"]}
        Weaknesses:
        {deck_data["swot_weaknesses"]}
        Opportunities:
        {deck_data["swot_opportunities"]}
        Threats:
        {deck_data["swot_threats"]}
"""
    # Slide 5
    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = \
        "Investment Thesis"

    slide.placeholders[1].text = \
        deck_data["investment_thesis"]

    # Slide 6
    slide = prs.slides.add_slide(
        prs.slide_layouts[1]
    )

    slide.shapes.title.text = \
        "Recommendation"

    slide.placeholders[1].text = f"""
        Recommendation:
        {deck_data['recommendation']}

        Risks:
        {deck_data['investment_risks']}
        """
    
    prs.save(output_file)

    return output_file
